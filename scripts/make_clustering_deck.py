#!/usr/bin/env python3
"""Generate Semantic Clustering presentation — PowerPoint-native design."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette (restrained) ──
NEAR_BLACK = RGBColor(0x1B, 0x1B, 0x1F)
DARK_GRAY = RGBColor(0x3B, 0x3B, 0x45)
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
FAINT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x3B, 0x5B, 0xDB)  # deep blue
ACCENT_LIGHT = RGBColor(0xDB, 0xE4, 0xFF)
TEAL = RGBColor(0x0D, 0x96, 0x88)
AMBER = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

MARGIN = Inches(0.9)
BODY_TOP = Inches(1.7)
BODY_W = SW - 2 * MARGIN


# ── helpers ──────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(sl, title, subtitle=None):
    """Consistent slide header: thin accent line + title."""
    # accent line
    line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(0.55), Inches(0.45), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.rotation = 0.0
    # title
    tb = sl.shapes.add_textbox(MARGIN + Inches(0.6), Inches(0.35), BODY_W - Inches(0.6), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NEAR_BLACK
    if subtitle:
        tb2 = sl.shapes.add_textbox(MARGIN + Inches(0.6), Inches(0.9), BODY_W - Inches(0.6), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MID_GRAY


def slide_num(sl, n, total=10):
    tb = sl.shapes.add_textbox(SW - Inches(1), SH - Inches(0.45), Inches(0.7), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{n}"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def textbox(sl, left, top, w, h):
    tb = sl.shapes.add_textbox(left, top, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_para(tf, text, size=14, color=DARK_GRAY, bold=False, space_after=Pt(6), first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = space_after
    return p


def add_rich(tf, parts, size=14, space_after=Pt(6)):
    """parts = list of (text, bold, color) tuples."""
    p = tf.add_paragraph()
    for text, bold, color in parts:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    p.space_after = space_after
    return p


def rounded_box(sl, left, top, w, h, fill=OFF_WHITE, border=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.text_frame.word_wrap = True
    return s


def image_callout(sl, left, top, w, h, label, filename):
    """Simple gray box with centered text — not a wireframe."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF3)
    s.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD8)
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # camera icon
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = MID_GRAY
    p.space_after = Pt(4)
    p2 = tf.add_paragraph()
    p2.text = filename
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER
    return s


# ═══════════════════════════════════════════════════════════════
# 1. TITLE
# ═══════════════════════════════════════════════════════════════
sl = blank()
slide_num(sl, 1)

# big accent block
block = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
block.fill.solid()
block.fill.fore_color.rgb = ACCENT
block.line.fill.background()

tf = textbox(sl, Inches(1.2), Inches(1.8), Inches(10), Inches(2))
add_para(tf, "Semantic Clustering for", size=38, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(2))
add_para(tf, "Multi-Cohort Data Harmonization", size=38, color=ACCENT, bold=True, space_after=Pt(20))
add_para(tf, "ddharmon  |  ARPA-H Activity 2", size=16, color=MID_GRAY)

# cohort names as simple text line
tf2 = textbox(sl, Inches(1.2), Inches(4.8), Inches(10), Inches(1))
add_para(tf2, "Arivale   ·   HPP   ·   UKBB   ·   TwinsUK   ·   All of Us", size=16, color=LIGHT_GRAY, first=True)
add_para(tf2, "5 cohorts  ·  questionnaires & demographics  ·  ~1,000+ fields combined", size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# 2. THE PROBLEM
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "The problem")
slide_num(sl, 2)

tf = textbox(sl, MARGIN, BODY_TOP, Inches(11), Inches(1))
add_para(tf, "Every cohort invented its own data dictionary. The same concept has different names,", size=18, color=DARK_GRAY, first=True, space_after=Pt(0))
add_para(tf, "different encodings, and different granularity. We need to find the common ground.", size=18, color=DARK_GRAY, space_after=Pt(16))

# concept divergence table — clean
headers = ["Concept", "Arivale", "UKBB", "HPP", "TwinsUK", "All of Us"]
rows = [
    ["Biological sex", "biological_sex", "sex", "sex", "sex", "Gender Identity"],
    ["Exercise", "exercise_frequency", "ipaq_total", "physical_activity_level", "exercise_frequency", "Physical Activity"],
    ["Smoking", "smoking_status", "current_tobacco_smoking", "smoking_habits", "smoking_history", "Smoking Status"],
    ["Education", "education_level", "qualifications", "(missing)", "education", "Education Level"],
]
tbl_shape = sl.shapes.add_table(len(rows) + 1, len(headers), MARGIN, Inches(3.0), Inches(11.5), Inches(2.0))
tbl = tbl_shape.table
for i, h in enumerate(headers):
    c = tbl.cell(0, i)
    c.text = h
    for p in c.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = MID_GRAY
    c.fill.solid()
    c.fill.fore_color.rgb = OFF_WHITE
for r, row in enumerate(rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r + 1, c_i)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = NEAR_BLACK if c_i == 0 else DARK_GRAY
            p.font.bold = c_i == 0
            if val == "(missing)":
                p.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
                p.font.italic = True

tf2 = textbox(sl, MARGIN, Inches(5.3), Inches(11), Inches(1.5))
add_para(tf2, "Clustering lets us discover these groupings automatically across all cohorts at once —", size=15, color=DARK_GRAY, first=True, space_after=Pt(0))
add_para(tf2, "no N*(N-1)/2 pairwise comparisons, and it surfaces coverage gaps immediately.", size=15, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# 3. PIPELINE AT A GLANCE
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Pipeline at a glance")
slide_num(sl, 3)

steps = [
    ("Load", "load_dictionary()", "Map CSV columns to\nsemantic roles"),
    ("Embed", "embed_dictionary()", "768-dim semantic vectors\nvia sentence-transformers"),
    ("Cluster", "cluster_dictionaries()", "Hierarchical agglomerative\n+ silhouette-driven cuts"),
    ("Explore", "Notebook", "UMAP, heatmaps,\ncluster inspection"),
]

box_w = Inches(2.5)
box_h = Inches(2.2)
gap = Inches(0.35)
total = int(box_w) * 4 + int(gap) * 3
start_x = (int(SW) - total) // 2

for i, (label, func, desc) in enumerate(steps):
    x = start_x + int((int(box_w) + int(gap)) * i)
    y = Inches(2.0)
    is_cluster = label == "Cluster"
    bx = rounded_box(sl, x, y, int(box_w), int(box_h),
                      fill=ACCENT_LIGHT if is_cluster else OFF_WHITE,
                      border=ACCENT if is_cluster else FAINT_GRAY)
    tf = bx.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # step number
    p = tf.paragraphs[0]
    p.text = f"Step {i+1}"
    p.font.size = Pt(10)
    p.font.color.rgb = ACCENT if is_cluster else LIGHT_GRAY
    p.font.bold = True
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT if is_cluster else NEAR_BLACK
    p2.space_after = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = func
    p3.font.size = Pt(11)
    p3.font.color.rgb = ACCENT if is_cluster else MID_GRAY
    p3.font.name = "Courier New"
    p3.space_after = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = desc
    p4.font.size = Pt(12)
    p4.font.color.rgb = DARK_GRAY

    # arrow between boxes
    if i < 3:
        ax = x + int(box_w) + int(gap * 0.15)
        arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   ax, int(y + int(box_h) // 2 - Inches(0.12)),
                                   int(gap * 0.7), Inches(0.24))
        arr.fill.solid()
        arr.fill.fore_color.rgb = FAINT_GRAY
        arr.line.fill.background()

# "today's focus" callout
tf3 = textbox(sl, start_x + int((int(box_w) + int(gap)) * 2), Inches(4.5), int(box_w), Inches(0.4))
add_para(tf3, "focus of this talk", size=11, color=ACCENT, bold=True, first=True)
tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

# bottom note
tf4 = textbox(sl, MARGIN, Inches(5.5), Inches(11), Inches(1.5))
add_para(tf4, "All four steps are implemented and validated. Phases 1-4 complete.", size=14, color=TEAL, bold=True, first=True)
add_para(tf4, "Clustering uses semantic vectors only — groups fields by what they measure, not how answers are encoded.", size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# 4. INSIDE cluster_dictionaries() — OVERVIEW
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Inside cluster_dictionaries()", "The orchestrator wires five substeps into a single call")
slide_num(sl, 4)

substeps = [
    ("Pool vectors", "Collect semantic vectors + FieldReference\nprovenance from all 5 cohorts into (N, 768) matrix"),
    ("Compute linkage", "Agglomerative hierarchical clustering\naverage linkage, cosine distance"),
    ("Suggest cuts", "Silhouette scan across 50 thresholds\nreturn top-3 scoring distances"),
    ("Extract clusters", "Flat cluster assignment at each cut\nwith per-cluster cohort coverage"),
    ("Label", "Word-frequency derived labels (offline)\nor optional LLM noun-phrase upgrade"),
]

for i, (name, desc) in enumerate(substeps):
    y = Inches(1.8) + Inches(1.05) * i
    # number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y, Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    cp = circ.text_frame.paragraphs[0]
    cp.text = str(i + 1)
    cp.font.size = Pt(14)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER

    # name
    tb = sl.shapes.add_textbox(MARGIN + Inches(0.6), y - Pt(2), Inches(2.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NEAR_BLACK

    # description
    tb2 = sl.shapes.add_textbox(Inches(3.6), y - Pt(2), Inches(4.5), Inches(0.9))
    tb2.text_frame.word_wrap = True
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY

    # connecting line (except last)
    if i < 4:
        ln = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  int(MARGIN) + Inches(0.18), int(y) + Inches(0.45),
                                  Pt(3), Inches(0.55))
        ln.fill.solid()
        ln.fill.fore_color.rgb = ACCENT_LIGHT
        ln.line.fill.background()

# right side: code snippet
code_box = rounded_box(sl, Inches(8.5), Inches(1.8), Inches(4.2), Inches(4.5), fill=OFF_WHITE, border=FAINT_GRAY)
ctf = code_box.text_frame
ctf.paragraphs[0].space_after = Pt(2)
p = ctf.paragraphs[0]
p.text = "Usage"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = LIGHT_GRAY

code = """hierarchy = cluster_dictionaries(
    embedded_dicts,
    linkage_method="average",
    llm_client=None,
    custom_cuts=[0.35, 0.60],
)

# Returns:
#   .linkage_matrix   (N-1, 4)
#   .cut_suggestions  top-3
#   .clusters_at_cuts {d: [...]}
#   .all_cohort_names"""

p2 = ctf.add_paragraph()
p2.text = code
p2.font.size = Pt(10)
p2.font.name = "Courier New"
p2.font.color.rgb = DARK_GRAY
p2.space_before = Pt(6)


# ═══════════════════════════════════════════════════════════════
# 5. SUBSTEP DEEP-DIVE: COMPUTE LINKAGE
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Substep: Compute linkage", "Agglomerative hierarchical clustering on pooled semantic vectors")
slide_num(sl, 5)

tf = textbox(sl, MARGIN, BODY_TOP, Inches(6), Inches(5))
add_para(tf, "How it works", size=18, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(12))

bullets = [
    [("All semantic vectors from all cohorts are stacked ", False, DARK_GRAY),
     ("into a single (N, 768) matrix.", False, DARK_GRAY)],
    [("Pairwise cosine distances ", False, DARK_GRAY),
     ("computed via scipy.spatial.distance.pdist().", False, DARK_GRAY)],
    [("Linkage matrix ", True, NEAR_BLACK),
     ("computed via scipy.cluster.hierarchy.linkage() with ", False, DARK_GRAY),
     ("average ", True, NEAR_BLACK),
     ("method.", False, DARK_GRAY)],
    [("Output: (N-1, 4) matrix ", False, DARK_GRAY),
     ("— each row records which two clusters merged, at what distance, and the size of the new cluster.", False, DARK_GRAY)],
]
for parts in bullets:
    add_rich(tf, parts, size=14, space_after=Pt(10))

add_para(tf, "Why average linkage + cosine?", size=16, color=NEAR_BLACK, bold=True, space_after=Pt(8))
add_rich(tf, [
    ("Average linkage ", True, NEAR_BLACK),
    ("is robust to outliers (vs single-linkage chaining). ", False, DARK_GRAY),
    ("Cosine distance ", True, NEAR_BLACK),
    ("matches how sentence-transformers measure semantic similarity — direction matters, not magnitude.", False, DARK_GRAY),
], size=14, space_after=Pt(10))

# method comparison
box = rounded_box(sl, Inches(7.5), BODY_TOP, Inches(5), Inches(2.5), fill=OFF_WHITE, border=FAINT_GRAY)
btf = box.text_frame
p = btf.paragraphs[0]
p.text = "Supported methods"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = MID_GRAY
p.space_after = Pt(10)

methods = [
    ("average (default)", "cosine", "Balanced — robust to noise"),
    ("complete", "cosine", "Conservative — tight clusters"),
    ("ward", "euclidean", "Minimizes variance — needs L2"),
]
for name, metric, note in methods:
    add_rich(btf, [
        (name, True, NEAR_BLACK),
        (f"  {metric}  ", False, LIGHT_GRAY),
        (note, False, DARK_GRAY),
    ], size=12, space_after=Pt(6))

# dendrogram note
note_box = rounded_box(sl, Inches(7.5), Inches(4.6), Inches(5), Inches(1.5), fill=RGBColor(0xFE, 0xF9, 0xC3), border=RGBColor(0xFD, 0xE6, 0x8A))
ntf = note_box.text_frame
p = ntf.paragraphs[0]
p.text = "Note on dendrograms"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = AMBER
p.space_after = Pt(4)
p2 = ntf.add_paragraph()
p2.text = "With 1,000+ fields, traditional dendrograms are unreadable. We skip them in favor of UMAP + flat cluster inspection at chosen distances."
p2.font.size = Pt(11)
p2.font.color.rgb = DARK_GRAY


# ═══════════════════════════════════════════════════════════════
# 6. SUBSTEP DEEP-DIVE: SILHOUETTE & CUT SELECTION
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Substep: Suggest cuts via silhouette analysis", "Data-driven threshold selection — no hardcoded distances")
slide_num(sl, 6)

# Left explanation
tf = textbox(sl, MARGIN, BODY_TOP, Inches(6.5), Inches(5))
add_para(tf, "Algorithm", size=18, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(10))

algo_steps = [
    [("Scan 50 evenly-spaced thresholds ", True, NEAR_BLACK),
     ("between min and max merge distances in the linkage matrix.", False, DARK_GRAY)],
    [("At each threshold, extract flat clusters ", False, DARK_GRAY),
     ("via fcluster(criterion='distance').", False, DARK_GRAY)],
    [("Filter extremes: ", True, NEAR_BLACK),
     ("skip if <2 clusters or >N/2 clusters (degenerate).", False, DARK_GRAY)],
    [("Compute silhouette score ", True, NEAR_BLACK),
     ("(sklearn, metric='cosine') — measures how well-separated clusters are. Range: -1 (bad) to +1 (crisp).", False, DARK_GRAY)],
    [("Return top 3 ", True, NEAR_BLACK),
     ("distances ranked by silhouette score as ", False, DARK_GRAY),
     ("CutSuggestion(distance, score, n_clusters).", False, MID_GRAY)],
]
for parts in algo_steps:
    add_rich(tf, parts, size=13, space_after=Pt(10))

add_para(tf, "The user then picks the cut that matches their analysis goal — or provides custom_cuts to override.", size=13, color=MID_GRAY, space_after=Pt(16))

# Right: image placeholder + interpretation guide
image_callout(sl, Inches(7.8), BODY_TOP, Inches(4.8), Inches(2.6),
              "Distance Sweep Table",
              "notebook cell 7 — distance vs cluster count vs universal clusters")

guide = rounded_box(sl, Inches(7.8), Inches(4.6), Inches(4.8), Inches(2.3), fill=OFF_WHITE, border=FAINT_GRAY)
gtf = guide.text_frame
p = gtf.paragraphs[0]
p.text = "Interpreting cut distances"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = NEAR_BLACK
p.space_after = Pt(8)

cuts = [
    ("d ~ 0.35  ", "Tight — near-identical fields only"),
    ("d ~ 0.50  ", "Balanced — related concepts grouped"),
    ("d ~ 0.70  ", "Loose — broad domain themes"),
]
for dist, interp in cuts:
    add_rich(gtf, [(dist, True, ACCENT), (interp, False, DARK_GRAY)], size=13, space_after=Pt(6))


# ═══════════════════════════════════════════════════════════════
# 7. SUBSTEP DEEP-DIVE: CLUSTER EXTRACTION & COVERAGE
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Substep: Extract clusters & track coverage", "Flat clusters at each cut distance, with cohort provenance")
slide_num(sl, 7)

tf = textbox(sl, MARGIN, BODY_TOP, Inches(6), Inches(4))
add_para(tf, "What extract_clusters() does", size=18, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(12))

extract_steps = [
    [("Calls fcluster(linkage_matrix, t=distance) ", False, DARK_GRAY),
     ("to assign every field to a flat cluster ID.", False, DARK_GRAY)],
    [("Groups FieldReferences by cluster. ", False, DARK_GRAY),
     ("Each FieldReference carries (cohort_name, variable_name, description).", False, DARK_GRAY)],
    [("Computes ", False, DARK_GRAY),
     ("cohort_coverage ", True, NEAR_BLACK),
     ("— dict mapping cohort name to member count per cluster.", False, DARK_GRAY)],
    [("Computes ", False, DARK_GRAY),
     ("missing_cohorts ", True, NEAR_BLACK),
     ("— which of the 5 cohorts have zero members in this cluster.", False, DARK_GRAY)],
]
for parts in extract_steps:
    add_rich(tf, parts, size=14, space_after=Pt(10))

# "Why this matters" box
why = rounded_box(sl, MARGIN, Inches(4.6), Inches(6), Inches(2.2), fill=ACCENT_LIGHT, border=ACCENT)
wtf = why.text_frame
p = wtf.paragraphs[0]
p.text = "Why coverage tracking matters"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = ACCENT
p.space_after = Pt(8)
add_rich(wtf, [
    ("Universal clusters ", True, ACCENT),
    ("(5/5 cohorts) = harmonizable concepts. ", False, DARK_GRAY),
    ("Partial clusters ", True, AMBER),
    ("= coverage gaps to flag. ", False, DARK_GRAY),
    ("Singleton clusters ", True, MID_GRAY),
    ("= cohort-specific fields with no equivalent.", False, DARK_GRAY),
], size=13, space_after=Pt(6))

# Right: example cluster
ex = rounded_box(sl, Inches(7.5), BODY_TOP, Inches(5), Inches(5.5), fill=OFF_WHITE, border=FAINT_GRAY)
etf = ex.text_frame
p = etf.paragraphs[0]
p.text = "Example: FieldCluster"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = MID_GRAY
p.space_after = Pt(10)

p2 = etf.add_paragraph()
p2.text = 'label = "Sex / Biological / Gender"'
p2.font.size = Pt(12)
p2.font.name = "Courier New"
p2.font.color.rgb = ACCENT
p2.space_after = Pt(6)

members = [
    "Arivale:  biological_sex",
    "HPP:      sex",
    "UKBB:     sex",
    "TwinsUK:  sex",
    "All of Us: Gender Identity",
]
for m in members:
    p3 = etf.add_paragraph()
    p3.text = f"  {m}"
    p3.font.size = Pt(11)
    p3.font.name = "Courier New"
    p3.font.color.rgb = DARK_GRAY
    p3.space_after = Pt(2)

p4 = etf.add_paragraph()
p4.space_before = Pt(10)
p4.text = "cohort_coverage:"
p4.font.size = Pt(11)
p4.font.bold = True
p4.font.color.rgb = MID_GRAY
p4.space_after = Pt(2)

p5 = etf.add_paragraph()
p5.text = "  Arivale: 1, HPP: 1, UKBB: 1,\n  TwinsUK: 1, All of Us: 1"
p5.font.size = Pt(11)
p5.font.name = "Courier New"
p5.font.color.rgb = TEAL
p5.space_after = Pt(6)

p6 = etf.add_paragraph()
p6.text = "missing_cohorts: []"
p6.font.size = Pt(11)
p6.font.name = "Courier New"
p6.font.color.rgb = TEAL


# ═══════════════════════════════════════════════════════════════
# 8. SUBSTEP DEEP-DIVE: LABELING
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Substep: Cluster labeling", "Two tiers — offline-first with optional LLM upgrade")
slide_num(sl, 8)

# Tier 1 box
t1 = rounded_box(sl, MARGIN, BODY_TOP, Inches(5.5), Inches(4.5), fill=OFF_WHITE, border=FAINT_GRAY)
t1tf = t1.text_frame

p = t1tf.paragraphs[0]
p.text = "Tier 1: Derived labels (default)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NEAR_BLACK
p.space_after = Pt(12)

derived_steps = [
    "Tokenize all member descriptions into words",
    "Remove 130+ stopwords (English + data dictionary terms like 'field', 'variable', 'participant')",
    "Count word frequency across all members",
    'Return top-3 words, title-cased, joined with " / "',
]
for s in derived_steps:
    add_rich(t1tf, [("  ", False, DARK_GRAY), (s, False, DARK_GRAY)], size=13, space_after=Pt(6))

p_ex = t1tf.add_paragraph()
p_ex.space_before = Pt(12)
p_ex.text = "Example"
p_ex.font.size = Pt(12)
p_ex.font.bold = True
p_ex.font.color.rgb = MID_GRAY
p_ex.space_after = Pt(4)

add_rich(t1tf, [
    ('Members: "Body mass index", "BMI measurement", "Body mass"', False, MID_GRAY),
], size=11, space_after=Pt(2))
add_rich(t1tf, [
    ("Output: ", False, MID_GRAY),
    ('"Body / Mass / Index"', True, ACCENT),
], size=13, space_after=Pt(8))

add_rich(t1tf, [
    ("No API key needed. ", True, TEAL),
    ("Runs instantly.", False, DARK_GRAY),
], size=13)

# Tier 2 box
t2 = rounded_box(sl, Inches(7), BODY_TOP, Inches(5.5), Inches(4.5), fill=OFF_WHITE, border=FAINT_GRAY)
t2tf = t2.text_frame

p = t2tf.paragraphs[0]
p.text = "Tier 2: LLM-upgraded labels (optional)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NEAR_BLACK
p.space_after = Pt(12)

llm_steps = [
    "For each cluster, assemble member names + descriptions",
    'Prompt: "Given these fields, provide a concise 2-4 word domain concept label (noun phrase)"',
    "Replace derived label with LLM response",
    "On failure, gracefully keep derived label",
]
for s in llm_steps:
    add_rich(t2tf, [("  ", False, DARK_GRAY), (s, False, DARK_GRAY)], size=13, space_after=Pt(6))

p_ex = t2tf.add_paragraph()
p_ex.space_before = Pt(12)
p_ex.text = "Example"
p_ex.font.size = Pt(12)
p_ex.font.bold = True
p_ex.font.color.rgb = MID_GRAY
p_ex.space_after = Pt(4)

add_rich(t2tf, [
    ("Same cluster as left", False, MID_GRAY),
], size=11, space_after=Pt(2))
add_rich(t2tf, [
    ("Output: ", False, MID_GRAY),
    ('"Body Mass Index"', True, ACCENT),
], size=13, space_after=Pt(8))

add_rich(t2tf, [
    ("Requires llm_client parameter. ", True, AMBER),
    ("Pass any BaseLLMClient.", False, DARK_GRAY),
], size=13)


# ═══════════════════════════════════════════════════════════════
# 9. UMAP + COVERAGE VISUALS
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "Visualization: UMAP + coverage heatmaps")
slide_num(sl, 9)

# Large UMAP placeholder
image_callout(sl, MARGIN, BODY_TOP, Inches(7), Inches(4.8),
              "UMAP Scatter Plot",
              "clustering_umap_dist{CUT_DISTANCE}.html\ncolor=cluster  shape=cohort  centroids labeled")

# Coverage heatmap
image_callout(sl, Inches(8.3), BODY_TOP, Inches(4.3), Inches(2.2),
              "Coverage Heatmap",
              "clustering_coverage_heatmap.png\nbinary presence + member counts")

# Size distribution
image_callout(sl, Inches(8.3), Inches(4.2), Inches(4.3), Inches(2.2),
              "Size & Breadth Distribution",
              "clustering_size_distribution.png\ncluster sizes + cohort breadth")

# Bottom note
tf = textbox(sl, MARGIN, Inches(6.7), Inches(7), Inches(0.5))
add_para(tf, "All generated in the validation notebook. UMAP is interactive (Plotly HTML). Heatmaps export as PNG.", size=12, color=LIGHT_GRAY, first=True)


# ═══════════════════════════════════════════════════════════════
# 10. WHAT THIS ENABLES
# ═══════════════════════════════════════════════════════════════
sl = blank()
title_bar(sl, "What this enables")
slide_num(sl, 10)

# Left: now
tf = textbox(sl, MARGIN, BODY_TOP, Inches(5.5), Inches(5))
add_para(tf, "Today", size=20, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(12))

now_items = [
    ("Coverage gap report", " — identify missing concepts per cohort before any manual mapping"),
    ("Cluster-guided pairwise", " — run detailed matching within clusters only, not exhaustive"),
    ("Exportable TSV", " — every field's cluster assignment for downstream tools"),
    ("Visual communication", " — UMAP + heatmaps make the harmonization landscape legible to reviewers"),
]
for bold_part, rest in now_items:
    add_rich(tf, [
        (bold_part, True, NEAR_BLACK),
        (rest, False, DARK_GRAY),
    ], size=14, space_after=Pt(10))

# Right: next
tf2 = textbox(sl, Inches(7), BODY_TOP, Inches(5.5), Inches(3.5))
add_para(tf2, "Next", size=20, color=NEAR_BLACK, bold=True, first=True, space_after=Pt(12))

next_items = [
    ("Value vector clustering", " — group by answer structure to inform transformations"),
    ("LLM label refinement", " — standardized biomedical concept names"),
    ("Cluster-to-pairwise bridge", " — auto-feed clusters into match_dictionaries()"),
    ("Expert-in-the-loop", " — review campaigns for edge cases"),
]
for bold_part, rest in next_items:
    add_rich(tf2, [
        (bold_part, True, NEAR_BLACK),
        (rest, False, DARK_GRAY),
    ], size=14, space_after=Pt(10))

# Status box
status = rounded_box(sl, Inches(7), Inches(5.5), Inches(5.5), Inches(1.2), fill=RGBColor(0xEC, 0xFD, 0xF5), border=TEAL)
stf = status.text_frame
p = stf.paragraphs[0]
p.text = "Phases 1-4 complete"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = TEAL
p.space_after = Pt(4)
p2 = stf.add_paragraph()
p2.text = "Load, embed, pairwise match, and cluster — all implemented and validated on 5 cohorts."
p2.font.size = Pt(13)
p2.font.color.rgb = DARK_GRAY


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output = "docs/semantic_clustering_deck.pptx"
prs.save(output)
print(f"Saved: {output}")
